import os
from pptx import Presentation
from pptx.util import Inches

# Path to the folder containing the images
image_folder = "E:/Games/Images/"

# Get the current working directory
script_dir = os.getcwd()

# Create a new PowerPoint presentation
presentation = Presentation()

# Get all files in the image folder
image_files = [file for file in os.listdir(image_folder) if os.path.isfile(os.path.join(image_folder, file))]

# Loop through the image files
for image_file in image_files:
    # Parse the filename to extract page and index numbers
    filename = os.path.splitext(image_file)[0]
    page_number, index_number = filename.split("_")[2], filename.split("_")[4]
    
    # Create a new slide
    slide_layout = presentation.slide_layouts[1]  # Choose the desired slide layout (index 1 is Title and Content)
    slide = presentation.slides.add_slide(slide_layout)
    
    # Add a title to the slide
    title = slide.shapes.title
    title.text = f"Image Page {page_number} Index {index_number}"
    
    # Add the image to the slide
    image_path = os.path.join(image_folder, image_file)
    slide.shapes.add_picture(image_path, left=Inches(1), top=Inches(1), width=Inches(6))  # Adjust the position and width as needed

# Save the PowerPoint presentation in the current working directory
output_file = os.path.join(script_dir, "output.pptx")
presentation.save(output_file)
