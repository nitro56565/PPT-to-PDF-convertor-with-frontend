import fitz
from pptx import Presentation
from pptx.util import Pt, Inches
import openai
import os
openai.api_key = "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"

def extract_important_points(text):
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=f"Summarize the following text in 4 points:\n{text}\n---\n",
        temperature=0.7,
        max_tokens=50,
        n=3,
        stop=None,
        frequency_penalty=0.5,
        presence_penalty=0.5
    )
    return [choice.text.strip() for choice in response.choices]

def extract_text_from_page(page):
    blocks = page.get_text("dict", flags=11)["blocks"]
    text = ""
    for block in blocks:
        for line in block["lines"]:
            for span in line["spans"]:
                text += span["text"]
        text += "\n"
    return text.strip()

def extract_slides_from_pdf(file_path, image_paths):
    try:
        with fitz.open(file_path) as pdf:
            slides = []
            for page_index, page in enumerate(pdf):
                text_blocks = page.get_text("dict", flags=11)["blocks"]
                if not text_blocks:  # Skip if no text blocks found on the page
                    continue
                
                title_block = max(
                    text_blocks,
                    key=lambda b: max([span["size"] for line in b["lines"] for span in line["spans"]])
                )
                title_text = " ".join([span["text"] for line in title_block["lines"] for span in line["spans"]])
                remaining_text = extract_text_from_page(page)
                important_points = extract_important_points(remaining_text)
                
                image_path = None
                if page_index < len(image_paths):
                    image_path = image_paths[page_index]
                
                slides.append((title_text, remaining_text, important_points, image_path))
            return slides
    except Exception as e:
        print(f"Error reading PDF file: {e}")
        return None

def extract_images_from_pdf(file_path, output_dir):
    images = []
    try:
        with fitz.open(file_path) as pdf:
            for page_index, page in enumerate(pdf):
                image_list = page.get_images()
                for image_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = pdf.extract_image(xref)
                    image_data = base_image["image"]
                    image_ext = base_image["ext"]
                    image_name = f"image_page_{page_index + 1}_index_{image_index + 1}.{image_ext}"
                    image_path = os.path.join(output_dir, image_name)
                    with open(image_path, "wb") as f:
                        f.write(image_data)
                    images.append(image_path)
        return images
    except Exception as e:
        print(f"Error extracting images from PDF: {e}")
        return None

def create_title_slide(prs, title):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

def create_content_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(16)
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    sentences = content.strip().split(".")
    for sentence in sentences[:4]:
        p = tf.add_paragraph()
        p.text = sentence.strip()
        p.font.size = Pt(12)
    if len(sentences) > 4:
        p = tf.add_paragraph()
        p.text = "[Continued on the next slide...]"
        p.level = 0
        p.font.size = Pt(12)

def create_presentation(pdf_file_path, ppt_file_path, image_output_dir):
    # Extract images from PDF and save them to the output directory
    image_paths = extract_images_from_pdf(pdf_file_path, image_output_dir)

    if image_paths is None:
        return

    # Clear existing slides and start with a blank presentation
    prs = Presentation()

    # Set the default theme to "Facet"
    prs.slide_layouts._default = prs.slide_layouts[6]

    # Extract slides from the PDF and associate images with their respective slides
    pdf_slides = extract_slides_from_pdf(pdf_file_path, image_paths)
    if pdf_slides is None:
        return

    # Create slides from the extracted PDF slides
    for index, (title, content, _, image_path) in enumerate(pdf_slides):
        if content.strip() != "":
            # Create a meaningful title phrase
            title_phrase = generate_meaningful_title(title)

            create_content_slide(prs, title_phrase, content)
        if image_path is not None:
            try:
                slide_layout = prs.slide_layouts[1]  # Choose the desired slide layout (index 1 is Title and Content)
                slide = prs.slides.add_slide(slide_layout)

                # Add a title to the slide
                image_title = os.path.splitext(os.path.basename(image_path))[0]
                # Create a meaningful title phrase for the image
                image_title_phrase = generate_meaningful_title(image_title)
                title = slide.shapes.title
                title.text = image_title_phrase

                # Add the image to the slide
                image_width = Inches(6)  # Adjust the desired width of the image
                image_height = Inches(4)  # Adjust the desired height of the image
                left = (prs.slide_width - image_width) / 2  # Calculate the left position to center the image
                top = (prs.slide_height - image_height) / 2  # Calculate the top position to center the image
                slide.shapes.add_picture(image_path, left=left, top=top, width=image_width, height=image_height)

                # Remove the image path from the list to prevent repetition
                image_paths.remove(image_path)
            except FileNotFoundError:
                print(f"Image file not found: {image_path}. Skipping to the next image.")

    # Save the PowerPoint presentation
    prs.save(ppt_file_path)
    print(f"PPT file saved to: {ppt_file_path}")

    # Extract images from PDF and save them to the output directory
    image_paths = extract_images_from_pdf(pdf_file_path, image_output_dir)


def generate_meaningful_title(original_title):
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=f"Generate a meaningful title phrase for:\n{original_title}\n---\n",
        temperature=0.7,
        max_tokens=20,
        n=1,
        stop=None,
        frequency_penalty=0.5,
        presence_penalty=0.5
    )
    return response.choices[0].text.strip()

    return original_title


pdf_file_path = "E:/Games/sample2.pdf"
ppt_file_path = "E:/Games/output.pptx"
image_output_dir = "E:/Games/Images/"

create_presentation(pdf_file_path, ppt_file_path, image_output_dir)