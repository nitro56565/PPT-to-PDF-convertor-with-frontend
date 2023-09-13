import sys
from PyQt5 import QtCore, QtGui, QtWidgets
import os
import fitz
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import openai

openai.api_key = "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"


def extract_important_points(text):
    response = openai.Completion.create(
        engine="text-curie-001",
        prompt=f"Summarize the following text in 4 points:\n{text}\n---\n",
        temperature=0.7,
        max_tokens=50,
        n=3,
        stop=None,
        frequency_penalty=0.5,
        presence_penalty=0.5
    )
    return [choice.text.strip() for choice in response.choices]


def extract_text_from_page(page):
    blocks = page.get_text("dict", flags=11)["blocks"]
    text = ""
    for block in blocks:
        for line in block["lines"]:
            for span in line["spans"]:
                text += span["text"]
        text += "\n"
    return text.strip()


def extract_slides_from_pdf(file_path, image_paths):
    try:
        with fitz.open(file_path) as pdf:
            slides = []
            for page_index, page in enumerate(pdf):
                text_blocks = page.get_text("dict", flags=11)["blocks"]
                if not text_blocks:  # Skip if no text blocks found on the page
                    continue

                title_block = max(
                    text_blocks,
                    key=lambda b: max([span["size"] for line in b["lines"] for span in line["spans"]])
                )
                title_text = " ".join([span["text"] for line in title_block["lines"] for span in line["spans"]])
                remaining_text = extract_text_from_page(page)
                important_points = extract_important_points(remaining_text)

                image_path = None
                if page_index < len(image_paths):
                    image_path = image_paths[page_index]

                slides.append((title_text, remaining_text, important_points, image_path))
            return slides
    except Exception as e:
        print(f"Error reading PDF file: {e}")
        return None


def extract_images_from_pdf(file_path, output_dir):
    images = []
    try:
        with fitz.open(file_path) as pdf:
            for page_index, page in enumerate(pdf):
                image_list = page.get_images()
                for image_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = pdf.extract_image(xref)
                    image_data = base_image["image"]
                    image_ext = base_image["ext"]
                    image_name = f"image_page_{page_index + 1}_index_{image_index + 1}.{image_ext}"
                    image_path = os.path.join(output_dir, image_name)
                    with open(image_path, "wb") as f:
                        f.write(image_data)
                    images.append(image_path)
        return images
    except Exception as e:
        print(f"Error extracting images from PDF: {e}")
        return None


def create_title_slide(prs, title):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title


def create_content_slide(prs, title, content, font_size, title_font_size, font_color, title_font_color):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(title_font_size)  # Set the title font size
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*title_font_color)  # Set the title font color
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    sentences = content.strip().split(".")
    for sentence in sentences[:4]:
        p = tf.add_paragraph()
        p.text = sentence.strip()
        p.font.size = Pt(font_size)  # Set the content font size
        p.font.color.rgb = RGBColor(*font_color)  # Set the content font color
    if len(sentences) > 4:
        p = tf.add_paragraph()
        p.text = "[Continued on the next slide...]"
        p.level = 0
        p.font.size = Pt(12)


def create_presentation(pdf_file_path, ppt_file_path, image_output_dir, font_size, title_font_size, font_color, title_font_color):
    # Extract images from PDF and save them to the output directory
    image_paths = extract_images_from_pdf(pdf_file_path, image_output_dir)

    if image_paths is None:
        return

    # Clear existing slides and start with a blank presentation
    prs = Presentation()

    # Set the default theme to "Facet"
    prs.slide_layouts._default = prs.slide_layouts[6]

    # Extract slides from the PDF and associate images with their respective slides
    pdf_slides = extract_slides_from_pdf(pdf_file_path, image_paths)
    if pdf_slides is None:
        return

    # Create slides from the extracted PDF slides
    for index, (title, content, _, image_path) in enumerate(pdf_slides):
        if content.strip() != "":
            # Create a meaningful title phrase
            title_phrase = generate_meaningful_title(title)

            create_content_slide(prs, title_phrase, content, font_size, title_font_size, font_color, title_font_color)
        if image_path is not None:
            try:
                slide_layout = prs.slide_layouts[1]  # Choose the desired slide layout (index 1 is Title and Content)
                slide = prs.slides.add_slide(slide_layout)

                # Add a title to the slide
                image_title = os.path.splitext(os.path.basename(image_path))[0]
                # Create a meaningful title phrase for the image
                image_title_phrase = generate_meaningful_title(image_title)
                title = slide.shapes.title
                title.text = image_title_phrase

                # Add the image to the slide
                image_width = Inches(6)  # Adjust the desired width of the image
                image_height = Inches(4)  # Adjust the desired height of the image
                left = (prs.slide_width - image_width) / 2  # Calculate the left position to center the image
                top = (prs.slide_height - image_height) / 2  # Calculate the top position to center the image
                slide.shapes.add_picture(image_path, left=left, top=top, width=image_width, height=image_height)

                # Remove the image path from the list to prevent repetition
                image_paths.remove(image_path)
            except FileNotFoundError:
                print(f"Image file not found: {image_path}. Skipping to the next image.")

    # Save the PowerPoint presentation
    prs.save(ppt_file_path)
    print(f"PPT file saved to: {ppt_file_path}")


def generate_meaningful_title(original_title):
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=f"Generate a meaningful title phrase for:\n{original_title}\n---\n",
        temperature=0.7,
        max_tokens=20,
        n=1,
        stop=None,
        frequency_penalty=0.5,
        presence_penalty=0.5
    )
    return response.choices[0].text.strip()


class Ui_MainWindow(object):
    def setupUi(self, MainWindow):
        MainWindow.setObjectName("MainWindow")
        MainWindow.resize(600, 500)
        self.centralwidget = QtWidgets.QWidget(MainWindow)
        self.centralwidget.setObjectName("centralwidget")
        self.frame = QtWidgets.QFrame(self.centralwidget)
        self.frame.setGeometry(QtCore.QRect(119, 99, 371, 281))
        self.frame.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.frame.setFrameShadow(QtWidgets.QFrame.Raised)
        self.frame.setObjectName("frame")

        # Set the background image using QPalette
        palette = QtGui.QPalette()
        background_image = QtGui.QPixmap("E:/Games/Frontend/sky.jpg")
        palette.setBrush(QtGui.QPalette.Background, QtGui.QBrush(background_image))
        MainWindow.setPalette(palette)

        self.BrowsePDF = QtWidgets.QPushButton(self.frame)
        self.BrowsePDF.setGeometry(QtCore.QRect(230, 30, 75, 23))
        self.BrowsePDF.setObjectName("BrowsePDF")
        self.Submit = QtWidgets.QPushButton(self.frame)
        self.Submit.setGeometry(QtCore.QRect(150, 220, 75, 23))
        self.Submit.setObjectName("Submit")
        self.FontSize = QtWidgets.QLabel(self.frame)
        self.FontSize.setGeometry(QtCore.QRect(100, 70, 51, 31))
        self.FontSize.setSizeIncrement(QtCore.QSize(2, 2))
        self.FontSize.setObjectName("FontSize")
        self.FontColour = QtWidgets.QLabel(self.frame)
        self.FontColour.setGeometry(QtCore.QRect(100, 100, 61, 31))
        self.FontColour.setObjectName("FontColour")
        self.TitleFontSize = QtWidgets.QLabel(self.frame)
        self.TitleFontSize.setGeometry(QtCore.QRect(100, 140, 71, 21))
        self.TitleFontSize.setObjectName("TitleFontSize")
        self.TitleFontColour = QtWidgets.QLabel(self.frame)
        self.TitleFontColour.setGeometry(QtCore.QRect(100, 170, 81, 16))
        self.TitleFontColour.setObjectName("TitleFontColour")
        self.FontSizenNo = QtWidgets.QComboBox(self.frame)  # Replace with QComboBox
        self.FontSizenNo.setGeometry(QtCore.QRect(200, 70, 71, 21))
        self.FontSizenNo.setObjectName("FontSizenNo")
        self.FontSizenNo.addItems(["2", "4", "6", "8", "10", "12", "14", "16", "18", "20", "22", "24", "26", "28", "30"])
        self.TitleFontColourNo = QtWidgets.QComboBox(self.frame)  # Replace with QComboBox
        self.TitleFontColourNo.setGeometry(QtCore.QRect(200, 170, 71, 21))
        self.TitleFontColourNo.setObjectName("TitleFontColourNo")
        self.TitleFontColourNo.addItems(["Red", "Blue", "Black", "Green", "Yellow", "White"])
        self.TitleFontSizeNo = QtWidgets.QComboBox(self.frame)  # Replace with QComboBox
        self.TitleFontSizeNo.setGeometry(QtCore.QRect(200, 130, 71, 21))
        self.TitleFontSizeNo.setObjectName("TitleFontSizeNo")
        self.TitleFontSizeNo.addItems(["2", "4", "6", "8", "10", "12", "14", "16", "18", "20", "22", "24", "26", "28", "30"])
        self.FontColourNo = QtWidgets.QComboBox(self.frame)  # Replace with QComboBox
        self.FontColourNo.setGeometry(QtCore.QRect(200, 100, 71, 21))
        self.FontColourNo.setObjectName("FontColourNo")
        self.FontColourNo.addItems(["Red", "Blue", "Black", "Green", "Yellow", "White"])
        self.FileName = QtWidgets.QLineEdit(self.frame)
        self.FileName.setGeometry(QtCore.QRect(70, 30, 151, 21))
        self.FileName.setObjectName("FileName")
        self.FileName.setReadOnly(True)  # Set the QLineEdit as read-only
        MainWindow.setCentralWidget(self.centralwidget)
        self.menubar = QtWidgets.QMenuBar(MainWindow)
        self.menubar.setGeometry(QtCore.QRect(0, 0, 800, 21))
        self.menubar.setObjectName("menubar")
        MainWindow.setMenuBar(self.menubar)
        self.statusbar = QtWidgets.QStatusBar(MainWindow)
        self.statusbar.setObjectName("statusbar")
        MainWindow.setStatusBar(self.statusbar)

        self.retranslateUi(MainWindow)
        QtCore.QMetaObject.connectSlotsByName(MainWindow)

        self.BrowsePDF.clicked.connect(self.browse_pdf)  # Connect the button click event to a custom method
        self.Submit.clicked.connect(self.submit_pdf)  # Connect the button click event to submit_pdf method

    def retranslateUi(self, MainWindow):
        _translate = QtCore.QCoreApplication.translate
        MainWindow.setWindowTitle(_translate("MainWindow", "MainWindow"))
        self.BrowsePDF.setText(_translate("MainWindow", "Browse PDF"))
        self.Submit.setText(_translate("MainWindow", "Submit"))
        self.FontSize.setText(_translate("MainWindow", "Font Size\n"
""))
        self.FontColour.setText(_translate("MainWindow", "Font Colour\n"
""))
        self.TitleFontSize.setText(_translate("MainWindow", "Title Font Size\n"
""))
        self.TitleFontColour.setText(_translate("MainWindow", "Title Font Colour"))

    def browse_pdf(self):
        file_dialog = QtWidgets.QFileDialog()
        file_dialog.setFileMode(QtWidgets.QFileDialog.ExistingFile)
        file_dialog.setNameFilter("PDF Files (*.pdf)")
        if file_dialog.exec_():
            selected_files = file_dialog.selectedFiles()
            if selected_files:
                selected_file = selected_files[0]
                file_name = QtCore.QFileInfo(selected_file).fileName()  # Extract the file name
                self.FileName.setText(file_name)
                self.pdf_file_path = selected_file

    def submit_pdf(self):
        ppt_file_path = "E:/Games/output.pptx"
        image_output_dir = "E:/Games/Images/"
        font_size = int(self.FontSizenNo.currentText())  # Get the selected font size
        title_font_size = int(self.TitleFontSizeNo.currentText())  # Get the selected title font size
        font_color = self.get_color_from_name(self.FontColourNo.currentText())  # Get the selected font color
        title_font_color = self.get_color_from_name(self.TitleFontColourNo.currentText())  # Get the selected title font color
        create_presentation(self.pdf_file_path, ppt_file_path, image_output_dir, font_size, title_font_size, font_color, title_font_color)
        print("Presentation created successfully.")

    def get_color_from_name(self, color_name):
        color_name = color_name.lower()
        if color_name == "red":
            return (255, 0, 0)
        elif color_name == "blue":
            return (0, 0, 255)
        elif color_name == "black":
            return (0, 0, 0)
        elif color_name == "green":
            return (0, 128, 0)
        elif color_name == "yellow":
            return (255, 255, 0)
        elif color_name == "white":
            return (255, 255, 255)
        else:
            return (0, 0, 0)  # Default to black if the color name is not recognized


class MainWindow(QtWidgets.QMainWindow):
    def __init__(self):
        super(MainWindow, self).__init__()

        # Create an instance of the generated UI
        self.ui = Ui_MainWindow()
        self.ui.setupUi(self)


if __name__ == "__main__":
    app = QtWidgets.QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())
